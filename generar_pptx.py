"""Genera presentacion.pptx — TP Minería de Datos CAECE"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# ── Paleta ──────────────────────────────────────────────────────
NAVY      = RGBColor(0x1e, 0x3a, 0x8a)
NAVY_DARK = RGBColor(0x0f, 0x24, 0x61)
TEAL      = RGBColor(0x0d, 0x94, 0x88)
TEAL_L    = RGBColor(0x5e, 0xea, 0xd4)
WHITE     = RGBColor(0xff, 0xff, 0xff)
GRAY_50   = RGBColor(0xf8, 0xfa, 0xfc)
GRAY_100  = RGBColor(0xf1, 0xf5, 0xf9)
GRAY_600  = RGBColor(0x47, 0x55, 0x69)
GRAY_800  = RGBColor(0x1e, 0x29, 0x3b)
AMBER     = RGBColor(0xd9, 0x77, 0x06)
GREEN     = RGBColor(0x16, 0xa3, 0x4a)
RED       = RGBColor(0xdc, 0x26, 0x26)

SW = Inches(13.33)
SH = Inches(7.5)
LOGO = "logo_caece.png"

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]


# ── Helpers ─────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def txbox(slide, text, l, t, w, h, size=16, bold=False,
          color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.name = "Calibri"; r.font.color.rgb = color
    return tb

def logo(slide, l=None, t=Inches(0.12), h=Inches(1.0)):
    if not os.path.exists(LOGO): return
    w = Inches(1.8)
    if l is None: l = SW - w - Inches(0.3)
    slide.shapes.add_picture(LOGO, l, t, w, h)

def header(slide, num_txt, title_txt):
    HH = Inches(1.35)
    rect(slide, 0, 0, SW, HH, NAVY)
    txbox(slide, num_txt,   Inches(0.5), Inches(0.12), Inches(9), Inches(0.4),
          size=11, bold=True, color=TEAL_L)
    txbox(slide, title_txt, Inches(0.5), Inches(0.48), Inches(10.2), Inches(0.78),
          size=26, bold=True, color=WHITE)
    rect(slide, 0, HH, SW, SH - HH, GRAY_50)
    logo(slide)
    return HH

def bullet_box(slide, items, l, t, w, h, bullet="→ ", size=15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = bullet + item
        r.font.size = Pt(size); r.font.name = "Calibri"
        r.font.color.rgb = GRAY_800

def card(slide, l, t, w, h, big_text, label, desc="", accent=NAVY):
    rect(slide, l, t, w, Inches(0.06), accent)
    rect(slide, l, t + Inches(0.06), w, h - Inches(0.06), WHITE)
    txbox(slide, big_text, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4),
          Inches(0.7), size=36, bold=True, color=accent)
    txbox(slide, label,    l + Inches(0.2), t + Inches(0.82), w - Inches(0.4),
          Inches(0.3), size=13, bold=True, color=GRAY_600)
    if desc:
        txbox(slide, desc, l + Inches(0.2), t + Inches(1.1), w - Inches(0.4),
              Inches(0.35), size=11, color=RGBColor(0x94, 0xa3, 0xb8))

def hbox(slide, l, t, w, h, title, body, fill=NAVY):
    rect(slide, l, t, w, h, fill)
    txbox(slide, title, l + Inches(0.3), t + Inches(0.18), w - Inches(0.6),
          Inches(0.4), size=16, bold=True, color=WHITE)
    txbox(slide, body,  l + Inches(0.3), t + Inches(0.55), w - Inches(0.6),
          h - Inches(0.7), size=14, color=WHITE, wrap=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, SW, SH, NAVY_DARK)
rect(sl, 0, SH - Inches(0.12), SW, Inches(0.12), TEAL)

if os.path.exists(LOGO):
    sl.shapes.add_picture(LOGO, SW/2 - Inches(1.5), Inches(0.6), Inches(3.0), Inches(1.4))

txbox(sl, "MINERÍA DE DATOS EMPRESARIALES · 2026",
      Inches(1), Inches(2.2), Inches(11.33), Inches(0.45),
      size=12, bold=True, color=TEAL_L, align=PP_ALIGN.CENTER)

txbox(sl, "¿A quién llamar?",
      Inches(0.5), Inches(2.7), Inches(12.33), Inches(1.4),
      size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(sl, "Un modelo predictivo para la captación de clientes bancarios",
      Inches(1.5), Inches(4.1), Inches(10.33), Inches(0.7),
      size=20, color=RGBColor(0xba, 0xc8, 0xff), align=PP_ALIGN.CENTER)

rect(sl, SW//2 - Inches(0.9), Inches(5.0), Inches(1.8), Inches(0.06), TEAL)

txbox(sl, "Lic. Lorena Lopez  ·  Lic. Gisela Poliak",
      Inches(1), Inches(5.2), Inches(11.33), Inches(0.45),
      size=16, color=RGBColor(0xba, 0xc8, 0xff), align=PP_ALIGN.CENTER)

txbox(sl, "Universidad CAECE  ·  2026  ·  1er Cuatrimestre",
      Inches(1), Inches(5.65), Inches(11.33), Inches(0.4),
      size=13, color=GRAY_600, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — EL DESAFÍO
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
HH = header(sl, "01 · El Desafío", "El banco desperdicia 88 de cada 100 llamadas")

CW = Inches(3.9); gap = Inches(0.28); top = HH + Inches(0.35)
card(sl, Inches(0.5),         top, CW, Inches(1.6), "49.732", "Clientes en el dataset",        "Registros históricos de campañas bancarias")
card(sl, Inches(0.5)+CW+gap,  top, CW, Inches(1.6), "~12%",   "Tasa de conversión actual",     "Solo 1 de cada 8 llamadas termina en venta", AMBER)
card(sl, Inches(0.5)+2*(CW+gap), top, CW, Inches(1.6), "~88%","Llamadas desperdiciadas",       "Tiempo y dinero perdidos en cada campaña",   RED)

hbox(sl, Inches(0.5), HH + Inches(2.1), SW - Inches(1.0), Inches(1.5),
     "🎯  La pregunta de este trabajo",
     "¿Podemos predecir a quién llamar y dejar de desperdiciar 88 llamadas de cada 100?\n"
     "Spoiler: sí podemos — y lo demostramos con datos reales.")


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — LOS DATOS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
HH = header(sl, "02 · Los Datos", "Conociendo el dataset y sus secretos")

# Mini stat cards (left column)
stats = [("49.732","Registros"),("15","Variables predictoras"),("12%","Clase positiva (sí)"),("88%","Clase negativa (no)")]
CW2 = Inches(2.7); RH = Inches(1.1)
for i,(val,lbl) in enumerate(stats):
    r,c = divmod(i,2)
    l = Inches(0.5) + c*(CW2+Inches(0.2))
    t = HH + Inches(0.3) + r*(RH+Inches(0.15))
    rect(sl, l, t, CW2, RH, WHITE)
    txbox(sl, val, l+Inches(0.15), t+Inches(0.1), CW2-Inches(0.3), Inches(0.55), size=28, bold=True, color=NAVY)
    txbox(sl, lbl, l+Inches(0.15), t+Inches(0.62), CW2-Inches(0.3), Inches(0.35), size=12, color=GRAY_600)

# Right column: hallazgos
RX = Inches(6.1); RW = SW - RX - Inches(0.4)
rect(sl, RX, HH+Inches(0.3), RW, SH-HH-Inches(0.5), WHITE)
txbox(sl, "Hallazgos en los datos", RX+Inches(0.25), HH+Inches(0.4), RW-Inches(0.5), Inches(0.4),
      size=16, bold=True, color=NAVY)
bullet_box(sl, [
    'Faltantes como "unknown" — no son nulos reales, son información',
    "Clase muy desbalanceada: 88% negativa — el Accuracy engaña",
    "Variables mixtas: numéricas (age, balance) y categóricas (job, education)",
    "Con datos desbalanceados, la métrica clave es AUC, no Accuracy",
], RX+Inches(0.2), HH+Inches(0.9), RW-Inches(0.4), Inches(2.4))


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — DATA LEAKAGE
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
HH = header(sl, "03 · La Trampa", "Data Leakage: el predictor perfecto... e inútil")

# Alert box
rect(sl, Inches(0.5), HH+Inches(0.3), SW-Inches(1.0), Inches(1.0), RGBColor(0xfe,0xf3,0xc7))
rect(sl, Inches(0.5), HH+Inches(0.3), Inches(0.08), Inches(1.0), AMBER)
txbox(sl, '⚠  Variable "duration": solo se conoce DESPUÉS de hacer la llamada.',
      Inches(0.7), HH+Inches(0.35), SW-Inches(1.4), Inches(0.4), size=15, bold=True, color=RGBColor(0x92,0x40,0x0e))
txbox(sl, "Usarla sería predecir el futuro con información del futuro. Decisión: eliminarla del modelo.",
      Inches(0.7), HH+Inches(0.72), SW-Inches(1.4), Inches(0.45), size=13, color=RGBColor(0x78,0x35,0x0f))

# Flow steps
steps = [("Dataset\noriginal","17 vars"),("❌ Eliminar\nduration","Data leakage"),
         ("División\n70 / 30","Stratify"),("One-Hot\nEncoding","Categóricas"),("✅ Modelo\nlisto","Honesto")]
FW = Inches(2.1); gap = Inches(0.08)
total_w = len(steps)*FW + (len(steps)-1)*gap
start_l = (SW - total_w) / 2
for i,(t1,t2) in enumerate(steps):
    l = start_l + i*(FW+gap)
    fill = RGBColor(0xfe,0xe2,0xe2) if i==1 else (RGBColor(0xdc,0xfc,0xe7) if i==4 else WHITE)
    rect(sl, l, HH+Inches(1.5), FW, Inches(0.9), fill)
    txbox(sl, t1, l+Inches(0.1), HH+Inches(1.55), FW-Inches(0.2), Inches(0.5),
          size=13, bold=True, color=(RED if i==1 else (GREEN if i==4 else NAVY)), align=PP_ALIGN.CENTER)
    txbox(sl, t2, l+Inches(0.1), HH+Inches(2.0), FW-Inches(0.2), Inches(0.3),
          size=11, color=GRAY_600, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        txbox(sl, "→", l+FW+Inches(0.01), HH+Inches(1.65), gap+Inches(0.05), Inches(0.4),
              size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Bottom cards
CW3 = Inches(3.9)
card(sl, Inches(0.5),         HH+Inches(2.6), CW3, Inches(1.35), "70%",  "Conjunto de entrenamiento", "Stratify para mantener proporción de clases", TEAL)
card(sl, Inches(0.5)+CW3+gap, HH+Inches(2.6), CW3, Inches(1.35), "30%",  "Conjunto de prueba",        "Nunca visto durante el entrenamiento")
card(sl, Inches(0.5)+2*(CW3+gap), HH+Inches(2.6), CW3, Inches(1.35), "+40", "Variables tras One-Hot Enc.","Sin orden artificial entre categorías", AMBER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — LOS MODELOS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
HH = header(sl, "04 · Los Modelos", "Dos modelos, una estrategia")

MCW = (SW - Inches(1.2)) / 2; MH = Inches(2.8)
for col,(ico,ttl,sub,items,params) in enumerate([
    ("🌳","CART — Árbol de decisión","Modelo controlado",
     ["Sin restricciones → 100% accuracy (overfitting): memorizó en vez de aprender",
      "Solución: limitar profundidad y exigir suficientes ejemplos por hoja",
      "class_weight='balanced' para atender la clase minoritaria"],
     ["max_depth=5","min_samples_leaf=20","class_weight='balanced'","random_state=42"]),
    ("🌲","Random Forest — Ensemble","100 árboles, mismos hiperparámetros",
     ["Consulta a 100 árboles y promedia sus votos — 100 segundas opiniones",
      "Reduce overfitting al promediar árboles con distintas submuestras",
      "Mismos hiperparámetros que CART para comparación justa"],
     ["n_estimators=100","max_depth=5","min_samples_leaf=20","class_weight='balanced'"])
]):
    l = Inches(0.5) + col*(MCW+Inches(0.2))
    rect(sl, l, HH+Inches(0.3), MCW, MH, WHITE)
    rect(sl, l, HH+Inches(0.3), MCW, Inches(0.08), NAVY if col==0 else TEAL)
    txbox(sl, ico+" "+ttl, l+Inches(0.2), HH+Inches(0.42), MCW-Inches(0.4), Inches(0.45),
          size=16, bold=True, color=NAVY)
    txbox(sl, sub, l+Inches(0.2), HH+Inches(0.82), MCW-Inches(0.4), Inches(0.3),
          size=12, color=GRAY_600)
    bullet_box(sl, items, l+Inches(0.15), HH+Inches(1.1), MCW-Inches(0.3), Inches(1.3), size=13)
    param_txt = "  ".join(params)
    txbox(sl, param_txt, l+Inches(0.15), HH+Inches(2.45), MCW-Inches(0.3), Inches(0.55),
          size=11, color=TEAL)

hbox(sl, Inches(0.5), HH+Inches(3.3), SW-Inches(1.0), Inches(1.0),
     "¿Por qué class_weight='balanced'?",
     "Con solo el 12% de positivos el modelo tiende a ignorar esa clase. El peso forzado hace que preste atención "
     "a los clientes que sí convierten — exactamente los que nos interesan.", TEAL)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — RESULTADOS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
HH = header(sl, "05 · Resultados", "¿Qué tan buenos son los modelos?")

# Table (left)
TW = Inches(6.0); TX = Inches(0.5); TY = HH+Inches(0.35)
rows = [("Métrica","🌳 CART","🌲 Random Forest",True),
        ("Accuracy","0.73","0.73",False),
        ("AUC-ROC","0.76","0.79  ✓ Mejor",False),
        ("Recall (clase sí)","~60%","~62%",False)]
RHT = Inches(0.52)
for i,(c1,c2,c3,hdr) in enumerate(rows):
    fill = NAVY if hdr else (GRAY_50 if i%2==0 else WHITE)
    col = WHITE if hdr else GRAY_800
    rect(sl, TX, TY+i*RHT, TW, RHT, fill)
    for j,(txt,w) in enumerate([(c1,Inches(2.4)),(c2,Inches(1.8)),(c3,Inches(1.8))]):
        lx = TX + (Inches(0)+Inches(2.4)*j if j==0 else (Inches(2.4)+(j-1)*Inches(1.8)))
        txbox(sl, txt, lx+Inches(0.1), TY+i*RHT+Inches(0.12), w-Inches(0.15), RHT-Inches(0.15),
              size=14 if not hdr else 13, bold=hdr or (i>0 and j==0),
              color=col if not (i==2 and j==2) else GREEN)

# Alert
rect(sl, TX, TY+4*RHT+Inches(0.2), TW, Inches(1.05), RGBColor(0xfe,0xf3,0xc7))
rect(sl, TX, TY+4*RHT+Inches(0.2), Inches(0.07), Inches(1.05), AMBER)
txbox(sl, "📊  ¿Por qué el Accuracy engaña?", TX+Inches(0.2), TY+4*RHT+Inches(0.27), TW-Inches(0.4), Inches(0.38),
      size=14, bold=True, color=RGBColor(0x92,0x40,0x0e))
txbox(sl, "Un modelo que siempre diga 'no convierte' tendría 88% de accuracy. El AUC mide la verdadera capacidad de discriminación.",
      TX+Inches(0.2), TY+4*RHT+Inches(0.6), TW-Inches(0.4), Inches(0.55), size=13, color=RGBColor(0x78,0x35,0x0f))

# Right column
RX2 = TX+TW+Inches(0.3); RW2 = SW-RX2-Inches(0.3)
rect(sl, RX2, HH+Inches(0.35), RW2, Inches(2.1), WHITE)
txbox(sl, "🏆  Random Forest gana en AUC", RX2+Inches(0.2), HH+Inches(0.45), RW2-Inches(0.4), Inches(0.42), size=15, bold=True, color=NAVY)
bullet_box(sl, [
    "AUC 0.79 vs 0.76 del CART — mejora consistente",
    "El ensemble promedia árboles distintos reduciendo el error por varianza",
    "Variable más importante: poutcome (campaña anterior)",
], RX2+Inches(0.15), HH+Inches(0.88), RW2-Inches(0.3), Inches(1.4), size=13)

hbox(sl, RX2, HH+Inches(2.6), RW2, Inches(1.1),
     "AUC 0.79 → usable en producción",
     "Por cada par de clientes (sí/no), el modelo ordena bien 79 de 100. Suficiente para priorizar llamadas en una campaña real.", GREEN)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — CONCLUSIONES
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
HH = header(sl, "06 · Impacto y Conclusiones", "Respuesta: sí podemos — y los números lo prueban")

# Lead text
txbox(sl, "Con 2.000 llamadas iguales:", Inches(0.5), HH+Inches(0.3), Inches(6.0), Inches(0.45),
      size=18, bold=True, color=NAVY)

# Impact cards
ICW = Inches(2.8); ICH = Inches(1.7)
for i,(lbl,val,acc) in enumerate([("📞  Llamando al azar","~234",GRAY_600),("🤖  Con el modelo","~805",GREEN)]):
    l = Inches(0.5) + i*(ICW+Inches(0.3))
    rect(sl, l, HH+Inches(0.85), ICW, Inches(0.08), acc)
    rect(sl, l, HH+Inches(0.93), ICW, ICH-Inches(0.08), WHITE)
    txbox(sl, lbl, l+Inches(0.15), HH+Inches(0.98), ICW-Inches(0.3), Inches(0.38), size=12, bold=True, color=GRAY_600)
    txbox(sl, val, l+Inches(0.15), HH+Inches(1.3), ICW-Inches(0.3), Inches(0.75), size=38, bold=True,
          color=(GREEN if i==1 else NAVY))
    txbox(sl, "conversiones", l+Inches(0.15), HH+Inches(2.05), ICW-Inches(0.3), Inches(0.4), size=12, color=GRAY_600)

hbox(sl, Inches(0.5), HH+Inches(2.75), Inches(6.0), Inches(1.1),
     "+244% de conversiones con el mismo presupuesto",
     "El modelo triplicó los resultados de la campaña sin agregar una sola llamada adicional.", GREEN)

# Right column
RX3 = Inches(6.9); RW3 = SW-RX3-Inches(0.4)
for i,(ico,ttl,items) in enumerate([
    ("✅","Lo que aprendimos",[
        "Overfitting y data leakage son los primeros enemigos a vencer",
        "Con datos desbalanceados: AUC > Accuracy como métrica clave",
        "Random Forest supera a CART de forma consistente",
    ]),
    ("🔮","Trabajo futuro",[
        "Probar XGBoost y LightGBM con validación cruzada",
        "Ajuste fino de hiperparámetros (GridSearchCV)",
        "Análisis de ROI por segmento de clientes",
    ]),
]):
    t = HH+Inches(0.3) + i*Inches(2.1)
    h = Inches(1.9)
    rect(sl, RX3, t, RW3, h, WHITE)
    txbox(sl, ico+"  "+ttl, RX3+Inches(0.2), t+Inches(0.15), RW3-Inches(0.4), Inches(0.4), size=15, bold=True, color=NAVY)
    bullet_box(sl, items, RX3+Inches(0.15), t+Inches(0.55), RW3-Inches(0.3), Inches(1.2), size=13)


# ── Guardar ──────────────────────────────────────────────────────
prs.save("presentacion.pptx")
print("✅  presentacion.pptx generado correctamente.")
